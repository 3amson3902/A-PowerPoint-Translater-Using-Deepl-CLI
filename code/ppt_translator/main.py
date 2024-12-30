"""
PowerPoint Translator
This script provides a GUI application to translate text within PowerPoint presentations using the DeepL translation service.
Modules:
    - pptx: Used to manipulate PowerPoint files.
    - deepl: Used to interact with the DeepL translation service.
    - os: Provides a way of using operating system dependent functionality.
    - tkinter: Provides a GUI toolkit for Python.
Functions:
    - initialize_deepl(input_lang, output_lang): Initializes the DeepLCLI with the specified input and output languages.
    - get_pptx_path(): Opens a file dialog to select a PowerPoint file and returns its path.
    - load_pptx(pptx_path): Loads a PowerPoint file from the given path and returns the presentation object.
    - read_translate_write_pptx(prs, deepl): Translates the text in the PowerPoint presentation using DeepL and returns the modified presentation.
    - save_pptx(prs, pptx_path, lang_var): Saves the modified PowerPoint presentation to a new file with a language-specific suffix.
    - check_lang_selected(): Checks if both input and output languages are selected and shows an error message if not.
    - main_run(): Main function to run the translation process, including file selection, translation, and saving the translated file.
GUI Elements:
    - root: The main window of the application.
    - inlang_var: A StringVar for the input language selection.
    - inlang_menu: A dropdown menu for selecting the input language.
    - outlang_var: A StringVar for the output language selection.
    - outlang_menu: A dropdown menu for selecting the output language.
    - open_button: A button to open the file dialog and start the translation process.
Usage:
    Run the script to open the GUI application. Select the input and output languages, import a PowerPoint file, and the application will translate the text within the presentation.
License:
    This script is provided "as is", without warranty of any kind, express or implied, including but not limited to the warranties of merchantability, fitness for a particular purpose, and noninfringement. In no event shall the authors be liable for any claim, damages, or other liability, whether in an action of contract, tort, or otherwise, arising from, out of, or in connection with the software or the use or other dealings in the software.
"""

from pptx import Presentation
from deepl import DeepLCLI
import os
import tkinter as tk
from tkinter import filedialog, messagebox

# Supported languages for DeepL
Lang = {
    "Japanese": "ja",
    "French": "fr",
    "Spanish": "es",
    "German": "de",
    "Chinese": "zh",
    "English": "en",
}


# Initialize DeepLCLI
def initialize_deepl(input_lang, output_lang):
    deepl = DeepLCLI(input_lang, output_lang)
    print(
        "DeepL Initialized" + input_lang + "To" + output_lang + deepl.translate("hello")
    )  # Test the translation
    return deepl


# Get the path to the PowerPoint file
def get_pptx_path():
    pptx_path = (
        filedialog.askopenfilename(  # Open a file dialog to select a PowerPoint file
            title="Select a PowerPoint file", filetypes=[("PowerPoint files", "*.pptx")]
        )
    )
    print("Path:" + pptx_path)  # Print the path to the selected file
    return pptx_path


# Load pptx
def load_pptx(pptx_path):  # pptx_path is the path to the PowerPoint file
    if not os.path.exists(pptx_path):  # Check if the file exists
        raise FileNotFoundError(
            f"The file {pptx_path} does not exist."
        )  # Raise an error if the file does not exist
    prs = Presentation(pptx_path)  # Load the PowerPoint file
    print("Presentation Loaded")  # Print a message
    return prs  # Return the presentation object


# Read, Translate and Write pptx
def read_translate_write_pptx(
    prs, deepl
):  # prs is the presentation object, deepl is the DeepLCLI object
    for slide in prs.slides:  # Iterate through all slides in the presentation
        for shape in slide.shapes:  # Iterate through all shapes in the slide
            if shape.has_text_frame:  # Check if the shape has a text frame
                for paragraph in (
                    shape.text_frame.paragraphs
                ):  # Iterate through all paragraphs in the text frame
                    for (
                        run
                    ) in paragraph.runs:  # Iterate through all runs in the paragraph
                        temp = deepl.translate(run.text)
                        print(run.text)  # Print the original text
                        run.text = temp  # Translate the text in the run
                        print(run.text)  # Print the translated text
    return prs  # Return the modified presentation


# Save pptx
def save_pptx(
    prs, pptx_path, lang_var
):  # prs is the presentation object, pptx_path is the path to the PowerPoint file, lang_var is the selected language
    new_pptx_path = (
        os.path.splitext(pptx_path)[0] + f"_{lang_var}.pptx"
    )  # Create a new file path with the selected language
    prs.save(new_pptx_path)  # Save the modified presentation to a new file
    return new_pptx_path  # Return the path to the new file


def check_lang_selected():  # Check if a language is selected
    if (
        inlang_var.get() == "Select Input Language"
    ):  # Check if the selected language is "Select Language"
        messagebox.showerror(
            "Error", "Please select a language."
        )  # Show an error message
    if (
        outlang_var.get() == "Select Output Language"
    ):  # Check if the selected language is "Select Language"
        messagebox.showerror(
            "Error", "Please select a language."
        )  # Show an error message
    return True


def main_run():
    try:
        pptx_path = get_pptx_path()  # Get the path to the PowerPoint file
        if check_lang_selected():
            prs = load_pptx(pptx_path)  # Load the PowerPoint file
            deepl = initialize_deepl(
                Lang[inlang_var.get()], Lang[outlang_var.get()]
            )  # Initialize DeepLCLI
            prs = read_translate_write_pptx(
                prs, deepl
            )  # Read, translate, and write the PowerPoint file
            new_pptx_path = save_pptx(
                prs, pptx_path, Lang[outlang_var.get()]
            )  # Save the translated PowerPoint file
            messagebox.showinfo(
                "Success", f"Translated file saved as {new_pptx_path}"
            )  # Show a success message
        else:
            messagebox.showinfo("Failed")  # Show a failed message
            return
    except Exception as e:
        messagebox.showerror("Error", f"An error occurred: {e}")


# Create the main window
root = tk.Tk()
root.title("PowerPoint Translator")

# Create a dropdown menu for input language selection
inlang_var = tk.StringVar(root)
inlang_var.set("Select Input Language")
inlang_menu = tk.OptionMenu(root, inlang_var, *Lang.keys())
inlang_menu.pack(pady=10)

# Create a dropdown menu for output language selection
outlang_var = tk.StringVar(root)
outlang_var.set("Select Output Language")
outlang_menu = tk.OptionMenu(root, outlang_var, *Lang.keys())
outlang_menu.pack(pady=10)

# Create a button to open the file dialog and run the main process
open_button = tk.Button(root, text="Import PowerPoint File", command=main_run)
open_button.pack(pady=20)

# Run the application
root.mainloop()
